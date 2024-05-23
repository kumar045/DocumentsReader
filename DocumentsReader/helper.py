import os
from pdf2image import convert_from_path
from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from openpyxl import load_workbook
from weasyprint import HTML

def get_text_size(text, font, max_width):
    lines = []
    words = text.split()
    while words:
        line = ''
        while words and font.getsize(line + words[0])[0] <= max_width:
            line = line + (words.pop(0) + ' ')
        lines.append(line)
    return lines

def create_image_with_text(text, font_path, max_width=800, padding=10):
    font = ImageFont.truetype(font_path, 15)
    lines = get_text_size(text, font, max_width - 2 * padding)
    max_height = sum(font.getsize(line)[1] for line in lines) + 2 * padding
    image = Image.new('RGB', (max_width, max_height), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    y = padding
    for line in lines:
        draw.text((padding, y), line, font=font, fill=(0, 0, 0))
        y += font.getsize(line)[1]
    return image

def convert_pdf_to_images(pdf_path, output_folder):
    images = convert_from_path(pdf_path)
    for i, image in enumerate(images):
        image_path = os.path.join(output_folder, f"page_{i+1}.png")
        image.save(image_path, "PNG")

def convert_docx_to_images(docx_path, output_folder):
    doc = Document(docx_path)
    font_path = "font/DejaVuSans.ttf"
    for i, para in enumerate(doc.paragraphs):
        image = create_image_with_text(para.text, font_path)
        image_path = os.path.join(output_folder, f"page_{i+1}.png")
        image.save(image_path, "PNG")


def convert_txt_to_images(txt_path, output_folder):
    with open(txt_path, "r") as file:
        text = file.read()
    font_path = "font/DejaVuSans.ttf"
    image = create_image_with_text(text, font_path)
    image_path = os.path.join(output_folder, "full_text.png")
    image.save(image_path, "PNG")

def convert_ppt_to_images(ppt_path, output_folder):
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        image = Image.new('RGB', (1280, 720), color = (255, 255, 255))
        draw = ImageDraw.Draw(image)
        left = top = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            draw.text((left, top), text, fill=(0, 0, 0))
            top += 40
        image_path = os.path.join(output_folder, f"slide_{i+1}.png")
        image.save(image_path, "PNG")

def convert_excel_to_images(excel_path, output_folder):
    workbook = load_workbook(excel_path)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        max_row = sheet.max_row
        max_col = sheet.max_column
        image = Image.new('RGB', (max_col*100, max_row*20), color = (255, 255, 255))
        draw = ImageDraw.Draw(image)
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            for col_idx, cell_value in enumerate(row, start=1):
                draw.text((col_idx*100, row_idx*20), str(cell_value), fill=(0, 0, 0))
        image_path = os.path.join(output_folder, f"{sheet_name}.png")
        image.save(image_path, "PNG")

def convert_html_to_image(html_path, output_folder):
    pdf_path = os.path.join(output_folder, "temp.pdf")
    HTML(html_path).write_pdf(pdf_path)
    convert_pdf_to_images(pdf_path, output_folder)
    os.remove(pdf_path)

def convert_file_to_images(file_path):
    output_folder = "output_folder"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    file_extension = file_path.split('.')[-1].lower()
    if file_extension == 'pdf':
        convert_pdf_to_images(file_path, output_folder)
    elif file_extension == 'docx':
        convert_docx_to_images(file_path, output_folder)
    elif file_extension == 'txt':
        convert_txt_to_images(file_path, output_folder)
    elif file_extension == 'ppt' or file_extension == 'pptx':
        convert_ppt_to_images(file_path, output_folder)
    elif file_extension == 'xlsx':
        convert_excel_to_images(file_path, output_folder)
    elif file_extension == 'html':
        convert_html_to_image(file_path, output_folder)
    else:
        raise ValueError("Unsupported file format")


